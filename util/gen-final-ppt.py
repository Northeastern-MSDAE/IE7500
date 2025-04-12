from pptx import Presentation
from pptx.util import Inches, Pt

# -----------------------------
# Presentation Data
# -----------------------------

presentation_data = [
    {
        "title": "Sentiment Analysis on Twitter Data",
        "content": ("IE7500 – Advanced Topics in Machine Learning\n"
                    "Subash Ramanathan, Ahantya Vempati, Yi Ren\n"
                    "April 2025\n"
                    "Code Repository: https://github.com/Northeastern-MSDAE/IE7500")
    },
    {
        "title": "Overview & Agenda",
        "content": ("Agenda:\n"
                    "• Introduction & Motivation\n"
                    "• Research & Influences\n"
                    "• Methodology & Implementation\n"
                    "• Results, Comparative Analysis & Error Analysis\n"
                    "• Conclusion & Future Work\n"
                    "• Q&A")
    },
    {
        "title": "Introduction & Motivation",
        "content": ("• Importance of Sentiment Analysis on Social Media\n"
                    "  - Understand public opinion from Twitter\n"
                    "  - Challenges: Informal language, slang, sarcasm, and short text\n"
                    "• Dataset: Sentiment140 (1.6 million tweets)\n"
                    "• Objectives: Robust preprocessing, effective feature extraction, diverse modeling")
    },
    {
        "title": "Research References & Influence",
        "content": ("Key Research Papers:\n"
                    "• Pak & Paroubek (2010): 'Twitter as a Corpus for Sentiment Analysis and Opinion Mining'\n"
                    "• Kouloumpis et al. (2011): 'Sentiment Analysis: The Good, the Bad and the OMG!'\n\n"
                    "Influence:\n"
                    "• Advanced tweet cleaning and normalization\n"
                    "• Enhanced feature extraction: n-grams, lexicon-based, microblogging-specific cues\n"
                    "• Considerations for error-aware training and data augmentation")
    },
    {
        "title": "Overall Pipeline Overview",
        "content": ("Pipeline Phases:\n"
                    "• Data Acquisition\n"
                    "• Preprocessing\n"
                    "• Feature Engineering\n"
                    "• Modeling\n"
                    "• Ensemble & Reporting\n\n"
                    "(Embedded Flowchart Image)")
    },
    {
        "title": "Data Preprocessing",
        "content": ("• Resources: NLTK (stopwords, wordnet)\n"
                    "• Cleaning Process:\n"
                    "  - Remove URLs, user mentions, retweet markers, hashtags\n"
                    "  - Convert emojis to text; lowercase; normalize elongated words\n"
                    "• Preprocessing Pipeline:\n"
                    "  - Tokenization, stop word removal, lemmatization\n\n"
                    "(Diagram/Code snippet illustration)")
    },
    {
        "title": "Exploratory Data Analysis (EDA)",
        "content": ("• Dataset Overview: Summary stats, class distribution\n"
                    "• Tweet Length Distribution:\n"
                    "  - High variability in tweet lengths\n"
                    "  - *Figure 1:* Histogram of tweet text lengths\n"
                    "• WordCloud Visualization:\n"
                    "  - Frequent tokens highlighted\n"
                    "  - *Figure 2:* WordCloud of tweet content")
    },
    {
        "title": "Feature Engineering",
        "content": ("• Baseline: TF-IDF vectorization (5,000 features)\n"
                    "• Custom Feature Transformers:\n"
                    "  - LexiconTransformer: Counts positive/negative words\n"
                    "  - MicrobloggingTransformer: Detects emoticons and abbreviations (e.g., 'OMG', 'BRB')\n"
                    "• Integration via FeatureUnion")
    },
    {
        "title": "Model Development",
        "content": ("• Traditional Models:\n"
                    "  - Logistic Regression, SVM, Random Forest, Naive Bayes\n"
                    "• Deep Learning Models:\n"
                    "  - LSTM, BiLSTM, CNN, Multi-Input, Attention-based models\n"
                    "• Model Tuning:\n"
                    "  - Hyperparameter tuning (GridSearchCV), Cross-validation\n\n"
                    "(Table summarizing performance)")
    },
    {
        "title": "Ensemble Methods",
        "content": ("• Strategy: Averaging predictions from traditional and deep learning models\n"
                    "• Benefits: Leverages complementary strengths, reduces variance, improves accuracy\n"
                    "• Results: Improved ROC & Precision-Recall metrics\n\n"
                    "(Figures: ROC Curve and Precision-Recall Curve)")
    },
    {
        "title": "Comparative Analysis",
        "content": ("• Traditional models: ~74% accuracy, efficient and stable\n"
                    "• Deep learning models: 70–72% accuracy, capture contextual nuances\n"
                    "• Insights:\n"
                    "  - Traditional methods are robust despite simpler architectures\n"
                    "  - Increased complexity yields marginal gains\n"
                    "  - Ensemble integration leverages complementary strengths")
    },
    {
        "title": "Error Analysis",
        "content": ("• Common Errors:\n"
                    "  - Ambiguous phrasing, sarcasm, informal language, lack of context\n"
                    "• Quantitative Insights:\n"
                    "  - Confusion matrix analysis indicates balanced errors\n"
                    "• Qualitative Examples:\n"
                    "  - Misclassified tweets due to sarcasm or mixed sentiment\n"
                    "• Improvement Strategies:\n"
                    "  - Enhanced features, contextual models (e.g., BERT), error-aware training\n\n"
                    "(Figure 5: Confusion Matrix)")
    },
    {
        "title": "Limitations and Future Work",
        "content": ("• Limitations:\n"
                    "  - Capturing nuances of informal language remains challenging\n"
                    "  - Trade-off between model complexity and efficiency\n"
                    "  - Binary classification might oversimplify sentiment\n"
                    "• Future Work:\n"
                    "  - Advanced preprocessing for slang, negation, and sarcasm\n"
                    "  - Incorporate transformer-based models (BERT, RoBERTa)\n"
                    "  - Deeper feature ablation, multi-class analysis, error-aware training")
    },
    {
        "title": "Conclusion and Final Remarks",
        "content": ("• Recap of contributions:\n"
                    "  - Robust pipeline from preprocessing to ensemble integration\n"
                    "  - Comparative analysis and error insights\n"
                    "• Final Thoughts:\n"
                    "  - Traditional models are strong; ensembles and advanced methods offer promise\n"
                    "\nThank you for your attention!")
    },
    {
        "title": "Q&A and Discussion",
        "content": "Questions & Answers"
    }
]

# -----------------------------
# Code to Generate the Presentation
# -----------------------------

def create_presentation(data, output_filename="Sentiment_Analysis_on_Twitter_Data.pptx"):
    prs = Presentation()
    slide_layout_title = prs.slide_layouts[0]  # Title Slide layout
    slide_layout_content = prs.slide_layouts[1]  # Title and Content layout

    for i, slide_data in enumerate(data):
        # Use Title Slide layout for the first slide, otherwise use content layout
        if i == 0:
            slide = prs.slides.add_slide(slide_layout_title)
            slide.shapes.title.text = slide_data["title"]
            # In Title Slide layout, the subtitle is usually placeholder with index 1
            subtitle = slide.placeholders[1]
            subtitle.text = slide_data["content"]
        else:
            slide = prs.slides.add_slide(slide_layout_content)
            slide.shapes.title.text = slide_data["title"]
            # Set content text into the second placeholder
            slide.shapes.placeholders[1].text = slide_data["content"]

    prs.save(output_filename)
    print(f"Presentation created: {output_filename}")

# Create the presentation using the defined data
if __name__ == "__main__":
    create_presentation(presentation_data)

